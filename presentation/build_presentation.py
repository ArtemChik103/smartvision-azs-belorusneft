"""
Build presentation pipeline:
1. Compile Markdown -> HTML via Marp CLI
2. Capture Pixel-Perfect 1280x720 (2x Retina) slides via Playwright Chromium
3. Assemble native 16:9 Microsoft PowerPoint presentation (*.pptx)
"""
import sys
import os
import shutil
import subprocess
import time
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

BASE_DIR = Path(__file__).resolve().parent.parent
PRES_DIR = BASE_DIR / "presentation"
SLIDES_MD = PRES_DIR / "slides.md"
SLIDES_HTML = PRES_DIR / "SmartVision-AZS-Presentation.html"
SLIDES_PPTX = PRES_DIR / "SmartVision-AZS-Presentation.pptx"
SLIDES_PDF = PRES_DIR / "SmartVision-AZS-Presentation.pdf"
SLIDES_IMG_DIR = PRES_DIR / "slide_images"


def build_presentation():
    PRES_DIR.mkdir(parents=True, exist_ok=True)
    if SLIDES_IMG_DIR.exists():
        shutil.rmtree(SLIDES_IMG_DIR)
    SLIDES_IMG_DIR.mkdir(parents=True, exist_ok=True)

    print("[1/3] Compiling Marp slides to HTML...")
    cmd = [
        "npx.cmd", "@marp-team/marp-cli",
        str(SLIDES_MD),
        "--no-stdin",
        "--html",
        "-o", str(SLIDES_HTML),
    ]
    res = subprocess.run(cmd, cwd=str(BASE_DIR), capture_output=True, text=True)
    if res.returncode != 0:
        print("[ERROR] Marp CLI failed:\n", res.stderr)
        return False

    print("[2/3] Capturing 2560x1440 (2x Retina) slides with Playwright...")
    from playwright.sync_api import sync_playwright

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        # Viewport matches Marp 16:9 native 1280x720 with 2x device scale = 2560x1440
        page = browser.new_page(viewport={"width": 1280, "height": 720}, device_scale_factor=2)
        
        file_url = f"file:///{str(SLIDES_HTML).replace(os.sep, '/')}"
        page.goto(file_url, wait_until="networkidle")
        
        # Hide Marp navigation overlay and progress bar
        page.add_style_tag(content="""
            .bespoke-marp-osc, .bespoke-progress-parent, .bespoke-marp-overview-header {
                display: none !important;
                opacity: 0 !important;
                visibility: hidden !important;
            }
            body, html {
                margin: 0 !important;
                padding: 0 !important;
                background-color: #0B1120 !important;
            }
        """)
        
        sections_count = page.evaluate("() => document.querySelectorAll('section').length")
        print(f"Found {sections_count} slides.")
        
        slide_images = []
        for idx in range(sections_count):
            img_path = SLIDES_IMG_DIR / f"slide_{idx+1:02d}.png"
            page.wait_for_timeout(150)
            page.screenshot(path=str(img_path))
            slide_images.append(img_path)
            print(f"  Rendered slide {idx+1}/{sections_count} -> {img_path.name}")
            page.keyboard.press("ArrowRight")
            
        browser.close()

    print("[3/3] Assembling 16:9 Microsoft PowerPoint presentation (*.pptx)...")
    prs = Presentation()
    prs.slide_width = Inches(13.3333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # Blank slide

    for img_path in slide_images:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(img_path),
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    prs.save(str(SLIDES_PPTX))
    
    # Export PDF via PIL
    from PIL import Image
    pil_imgs = [Image.open(p).convert("RGB") for p in slide_images]
    if pil_imgs:
        pil_imgs[0].save(str(SLIDES_PDF), save_all=True, append_images=pil_imgs[1:])
    
    print(f"\n[SUCCESS] Presentation build complete!")
    print(f"  PPTX: {SLIDES_PPTX} ({round(SLIDES_PPTX.stat().st_size / (1024*1024), 2)} MB)")
    print(f"  PDF:  {SLIDES_PDF} ({round(SLIDES_PDF.stat().st_size / (1024*1024), 2)} MB)")
    print(f"  HTML: {SLIDES_HTML}")
    return True


if __name__ == "__main__":
    build_presentation()
